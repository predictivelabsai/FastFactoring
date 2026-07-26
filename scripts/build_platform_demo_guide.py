"""Build the English Factor AI Platform Demo guide from Playwright captures."""

from __future__ import annotations

from datetime import date
from pathlib import Path

import markdown
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from weasyprint import HTML

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
IMG_DIR = DOCS / "img" / "platform-demo"
DATE = date.today().isoformat()
BASE = DOCS / f"factor_ai_platform_demo_user_guide_{DATE}_en"
TITLE = "Factor AI Platform Demo — User Guide"
GREEN = RGBColor(0x1F, 0x5D, 0x43)
INK = RGBColor(0x14, 0x23, 0x1B)
MUTED = RGBColor(0x5C, 0x68, 0x60)
PARCH = RGBColor(0xF7, 0xF6, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SECTIONS = [
    ("1. Borrower / Supplier", [
        ("01-supplier-ai-upload.png", "Start in Factorio AI",
         ["The Supplier demo opens in one AI workspace—there is no marketplace or triage queue.",
          "Select the paperclip and upload a digital PDF, image, JSON, or text invoice.",
          "The demo uses synthetic invoices; digital PDFs are converted to Markdown before xAI extraction."]),
        ("02-supplier-ai-offer.png", "Upload an invoice and receive an offer",
         ["Factorio AI extracts supplier, debtor, invoice number, amount, dates, bank details, and payment terms.",
          "The response clearly says “You’re pre-approved!” and presents an indicative offer table.",
          "Download PDF, Accept, and Change remain visible; final terms are subject to verification."]),
        ("03-supplier-change-terms.png", "Change the amount or financing period",
         ["Open Change to request a lower amount paid today or a different period.",
          "The indicative fee updates from the requested amount and days.",
          "Add contact details and confirm registry information before acceptance."]),
        ("04-supplier-profile-integrations.png", "Profile, banking, and accounting connections",
         ["My Profile contains a hypothetical company identity and masked settlement bank information.",
          "Lloyds Bank is labelled as a planned Open Banking connection.",
          "QuickBooks, Xero, and Sage are clearly labelled hypothetical demo integrations."]),
        ("05-supplier-applications.png", "Track applications and contracts",
         ["My Applications shows synthetic financing requests created from accepted invoices.",
          "Each accepted request can expose a one-page financing contract between Factorio Ltd and the extracted supplier.",
          "Downloaded PDFs use supplier, Factorio reference, date, and document-type slugs."]),
    ]),
    ("2. Investor", [
        ("06-investor-ai-performance.png", "Ask Investor AI about performance",
         ["Investor AI is grounded only in the selected demo investor’s positions and computed metrics.",
          "Eval question shown: “How are my investments performing?”",
          "Expected and realised returns are distinguished; the assistant does not guarantee outcomes."]),
        ("07-investor-ai-concentration.png", "Ask about concentration and risk",
         ["Eval question shown: “Where is my portfolio most concentrated?”",
          "Debtor and sector exposure percentages are calculated deterministically before reaching the model.",
          "Copy and Share create portable session transcripts; shared links are read-only snapshots."]),
        ("08-investor-portfolio.png", "Review the portfolio cockpit",
         ["See account value, net annual return, active invested, expected outstanding, and realised results.",
          "Aging, payment habits, grades, due dates, and position status support deeper AI questions.",
          "Use the demo investor selector to compare isolated synthetic portfolios."]),
        ("09-investor-marketplace.png", "Inspect available invoices",
         ["Browse synthetic open invoices by debtor, sector, risk grade, term, and estimated return.",
          "Funding progress and invoice economics support manual selection.",
          "Returns are estimates and invoice financing retains credit, fraud, dilution, concentration, and liquidity risk."]),
        ("10-investor-autoinvest.png", "Configure risk-aware Auto-invest",
         ["Choose conservative, balanced, or growth preferences plus grade, term, return, sector, debtor-concentration, and per-invoice limits.",
          "Investor AI ranks only eligible invoices and explains each proposed allocation.",
          "The preview does not place money; review or change preferences before execution."]),
        ("11-investor-statement.png", "Reconcile investment activity",
         ["The statement records investments and settlements with invoice and counterparty references.",
          "Filter by date, type, counterparty, invoice, and amount.",
          "Export CSV for external analysis or reconciliation."]),
    ]),
    ("3. Admin", [
        ("12-admin-console.png", "Operate the back office",
         ["Admin sees platform KPIs, queues, exceptions, and operational navigation in one workspace.",
          "Navigation groups are collapsible and can be minimized to keep a complex toolset manageable.",
          "Global settings include the USD, GBP, EUR, or UZS display-currency override."]),
        ("13-admin-agent-fleet.png", "Coordinate the Agent Fleet",
         ["Fourteen specialists span orchestration, origination, decisioning, servicing, oversight, and growth.",
          "Agents use grounded tools and preserve approval gates for credit, communication, publishing, spend, and money movement.",
          "The activity feed makes agent operations auditable."]),
        ("14-admin-seo-agent.png", "Run marketing and SEO agents",
         ["Browser-driven eval question shown: “Build a supplier invoice-finance keyword cluster.”",
          "SEO & AI Search and Paid Marketing agents create research and drafts without silently publishing or activating spend.",
          "Campaign work should include audience, intent, evidence, measurement, guardrails, and explicit approval."]),
        ("15-admin-skills-editor.png", "Inspect and edit agent skills",
         ["Open Agent Skills to view each specialist’s current Markdown instructions.",
          "Admin can save a database-backed prompt version and review change history.",
          "Revert restores a selected earlier version while preserving an audit trail."]),
        ("16-admin-processing.png", "Process and approve invoices",
         ["The processing queue supports verification, exceptions, risk review, and funding readiness.",
          "AI may recommend or draft; final financial and credit actions remain approval-gated.",
          "Synthetic data makes the workflow safe to explore in the platform demo."]),
        ("17-admin-accounting.png", "Reconcile accounting activity",
         ["Accounting tools cover journal review, reconciliation, receivables, fees, reserves, and settlement evidence.",
          "Agents must not create balancing plugs or invent transactions.",
          "Use source records and approval history to resolve differences."]),
        ("18-admin-integrations.png", "Manage platform integrations",
         ["Integration cards show demo connectivity for accounting, banking, communications, and document services.",
          "Connection status is demonstrative and must not be treated as a live third-party authorization.",
          "The invoice ingestion endpoint supports structured external demo data."]),
        ("19-admin-audit.png", "Review the audit trail",
         ["Audit records identify actor, action, entity, detail, and time.",
          "Use the log to inspect human and agent activity, approvals, and configuration changes.",
          "Prompt-version history and action logs support governance of the agentic platform."]),
    ]),
]


def markdown_text() -> str:
    parts = [
        f"# {TITLE}",
        "",
        f"**Generated {DATE} · Live demo: https://factorio.co.uk**",
        "",
        "> Demo notice: every company, invoice, offer, integration, investment, and AI response shown here is synthetic or illustrative. "
        "The platform does not guarantee funding, payment timing, or investment returns.",
        "",
        "This guide is a browser-driven tour of the Factor AI Platform demo. It follows the three primary roles and uses reviewed evaluation questions to demonstrate AI behavior.",
        "",
        "![Factor AI Platform demo](img/platform-demo/00-platform-demo.png)",
        "",
    ]
    for section, slides in SECTIONS:
        parts += ["---", "", f"# {section}", ""]
        for image, title, bullets in slides:
            parts += [
                f"## {title}", "",
                f"![{title}](img/platform-demo/{image})", "",
                *[f"- {item}" for item in bullets], "",
            ]
    parts += [
        "---", "", "# Suggested demo script", "",
        "1. Sign in as **Supplier**, upload a sample invoice, inspect the offer, and open My Profile.",
        "2. Switch to **Investor**, ask the two evaluation questions shown above, then compare Portfolio and Auto-invest.",
        "3. Switch to **Admin**, open Agent Fleet, run a specialist prompt, inspect Agent Skills, then review Processing, Accounting, Integrations, and Audit.",
        "",
        "The demo credentials and one-click roles are available from the Sign in page. Use only synthetic information.",
    ]
    return "\n".join(parts)


def html_text(source: str) -> str:
    body = markdown.markdown(source, extensions=["tables", "sane_lists"])
    css = (DOCS / "assets" / "guide.css").read_text(encoding="utf-8")
    return (
        "<!doctype html><html lang='en'><head><meta charset='utf-8'>"
        f"<title>{TITLE}</title><style>{css}</style></head><body>{body}</body></html>"
    )


def _bg(slide, prs, color=PARCH):
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def _text(slide, text, left, top, width, height, size, color=INK, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return frame


def _add_contained(slide, path: Path, left, top, width, height):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    if ratio > box_ratio:
        draw_w, draw_h = width, width / ratio
    else:
        draw_h, draw_w = height, height * ratio
    slide.shapes.add_picture(
        str(path), left + (width - draw_w) / 2, top + (height - draw_h) / 2,
        width=draw_w, height=draw_h,
    )


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    _bg(slide, prs, GREEN)
    _text(slide, "FACTOR AI PLATFORM DEMO", Inches(.8), Inches(1.8), Inches(11.8),
          Inches(.5), 14, WHITE, True)
    _text(slide, "User Guide", Inches(.8), Inches(2.45), Inches(11.8),
          Inches(1), 42, WHITE, True)
    _text(slide, "Supplier · Investor · Admin", Inches(.8), Inches(3.55),
          Inches(11.8), Inches(.6), 24, WHITE)
    _text(slide, f"Synthetic demonstration data · {DATE}", Inches(.8), Inches(5.8),
          Inches(11.8), Inches(.5), 13, WHITE)

    for section, slides in SECTIONS:
        divider = prs.slides.add_slide(blank)
        _bg(divider, prs, GREEN)
        _text(divider, section, Inches(.9), Inches(2.65), Inches(11.5),
              Inches(1.2), 38, WHITE, True)
        for image, title, bullets in slides:
            slide = prs.slides.add_slide(blank)
            _bg(slide, prs)
            _text(slide, section.upper(), Inches(.55), Inches(.25), Inches(12.2),
                  Inches(.3), 10, GREEN, True)
            _text(slide, title, Inches(.55), Inches(.62), Inches(12.2),
                  Inches(.65), 24, GREEN, True)
            _add_contained(slide, IMG_DIR / image, Inches(.55), Inches(1.4),
                           Inches(7.4), Inches(5.55))
            frame = _text(slide, "• " + bullets[0], Inches(8.25), Inches(1.55),
                          Inches(4.45), Inches(4.9), 14)
            for bullet in bullets[1:]:
                paragraph = frame.add_paragraph()
                paragraph.text = "• " + bullet
                paragraph.space_before = Pt(13)
                paragraph.runs[0].font.size = Pt(14)
                paragraph.runs[0].font.color.rgb = INK
            _text(slide, "DEMO · SYNTHETIC DATA", Inches(8.25), Inches(6.55),
                  Inches(4.4), Inches(.3), 9, MUTED, True)
    prs.save(path)


def main() -> None:
    source = markdown_text()
    BASE.with_suffix(".md").write_text(source, encoding="utf-8")
    rendered = html_text(source)
    BASE.with_suffix(".html").write_text(rendered, encoding="utf-8")
    HTML(string=rendered, base_url=str(DOCS)).write_pdf(BASE.with_suffix(".pdf"))
    build_pptx(BASE.with_suffix(".pptx"))
    for extension in ("md", "html", "pdf", "pptx"):
        path = BASE.with_suffix("." + extension)
        print(f"{path.name}: {path.stat().st_size / 1024:.0f} KB")


if __name__ == "__main__":
    main()
